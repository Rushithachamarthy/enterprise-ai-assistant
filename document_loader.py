# document_loader.py
# This file reads different types of files and extracts text
# Handles images with OCR (now works on Streamlit Cloud too)

import pandas as pd
from pypdf import PdfReader
from docx import Document
from pptx import Presentation
from PIL import Image
import pytesseract
import io

# REMOVED THE WINDOWS-SPECIFIC LINE:
# pytesseract.pytesseract.tesseract_cmd = r'C:\Program Files\Tesseract-OCR\tesseract.exe'
# → Not needed anymore. Streamlit Cloud installs Tesseract via packages.txt

def load_document(file):
    file_name = file.name.lower()
    metadata = {"file_type": file_name.split('.')[-1]}

    if file_name.endswith(".pdf"):
        reader = PdfReader(file)
        text = "\n".join(page.extract_text() or "" for page in reader.pages)
        metadata["page_count"] = len(reader.pages)
        return text, metadata

    elif file_name.endswith(".docx"):
        doc = Document(file)
        paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
        text = "\n".join(paragraphs)
        metadata["page_count"] = len(paragraphs)
        return text, metadata

    elif file_name.endswith(".pptx"):
        prs = Presentation(file)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        metadata["slide_count"] = len(prs.slides)
        return "\n".join(text), metadata

    elif file_name.endswith(".csv"):
        df = pd.read_csv(file)
        text = df.to_string()
        metadata["row_count"] = len(df)
        metadata["column_count"] = len(df.columns)
        return text, metadata

    elif file_name.endswith(".xlsx"):
        df = pd.read_excel(file)
        text = df.to_string()
        metadata["row_count"] = len(df)
        metadata["column_count"] = len(df.columns)
        return text, metadata

    elif file_name.endswith(".txt"):
        text = file.read().decode("utf-8")
        metadata["line_count"] = len(text.splitlines())
        return text, metadata

    elif file_name.endswith((".jpg", ".jpeg", ".png")):
        file.seek(0)
        image = Image.open(io.BytesIO(file.read()))
        text = pytesseract.image_to_string(image)
        metadata["is_image"] = True  # Special flag for images
        return text, metadata

    return "Could not extract text from this file.", metadata